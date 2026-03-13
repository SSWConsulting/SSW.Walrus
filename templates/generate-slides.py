#!/usr/bin/env python3
"""Generate an SSW-branded PPTX slide deck from consolidated survey analysis JSON.

Usage:
    python3 templates/generate-slides.py <consolidated.json> <output.pptx>

Example:
    python3 templates/generate-slides.py \
        surveys/engagement/2026-03-04/analysis/consolidated.json \
        surveys/engagement/2026-03-04/dashboard/engagement.pptx
"""

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# SSW Brand Colors
SSW_RED = RGBColor(0xCC, 0x41, 0x41)
SSW_CHARCOAL = RGBColor(0x33, 0x33, 0x33)
SSW_GRAY = RGBColor(0x66, 0x66, 0x66)
SSW_LIGHT_GRAY = RGBColor(0xF4, 0xF4, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SSW_GREEN = RGBColor(0x05, 0x96, 0x69)
SSW_AMBER = RGBColor(0xF5, 0x9E, 0x0B)
SSW_RED_50 = RGBColor(0xFC, 0xE9, 0xE9)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
FONT_NAME = "Calibri"


# ---------------------------------------------------------------------------
# Helpers for resolving nested JSON paths
# ---------------------------------------------------------------------------

def resolve(data, *keys, default=None):
    """Walk nested dicts/lists to find a value. Tries multiple key paths."""
    for key in keys:
        if isinstance(key, str) and "." in key:
            parts = key.split(".")
            val = data
            for p in parts:
                if isinstance(val, dict):
                    val = val.get(p)
                else:
                    val = None
                    break
            if val is not None:
                return val
        elif isinstance(data, dict):
            val = data.get(key)
            if val is not None:
                return val
    return default


# ---------------------------------------------------------------------------
# Text helpers
# ---------------------------------------------------------------------------

def set_text(textframe, text, font_size=14, bold=False, color=SSW_CHARCOAL, alignment=PP_ALIGN.LEFT):
    """Set text in a textframe with formatting."""
    textframe.clear()
    p = textframe.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT_NAME
    return p


def add_paragraph(textframe, text, font_size=14, color=SSW_CHARCOAL, bold=False,
                   space_before=4, space_after=2):
    """Add a formatted paragraph."""
    p = textframe.add_paragraph()
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT_NAME
    return p


def write_bullets(textframe, items, font_size=14, color=SSW_CHARCOAL, bold=False, prefix="• "):
    """Write a list of bullet strings into a textframe, replacing existing content."""
    for i, item in enumerate(items):
        text = f"{prefix}{item}"
        if i == 0:
            textframe.paragraphs[0].text = ""
            run = textframe.paragraphs[0].add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = FONT_NAME
        else:
            add_paragraph(textframe, text, font_size=font_size, color=color, bold=bold)


# ---------------------------------------------------------------------------
# Slide furniture
# ---------------------------------------------------------------------------

def add_red_bar(slide, top=Inches(0)):
    """Add SSW red accent bar at the top of a slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), top, SLIDE_WIDTH, Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = SSW_RED
    shape.line.fill.background()


def add_slide_number(slide, number):
    """Add slide number to bottom-right."""
    txBox = slide.shapes.add_textbox(
        Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(number)
    run.font.size = Pt(10)
    run.font.color.rgb = SSW_GRAY
    run.font.name = FONT_NAME


def add_section_title(slide, title, subtitle=None):
    """Add a section title with SSW styling."""
    add_red_bar(slide)

    txBox = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8)
    )
    tf = txBox.text_frame
    set_text(tf, title, font_size=28, bold=True, color=SSW_CHARCOAL)

    if subtitle:
        txBox2 = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.4)
        )
        tf2 = txBox2.text_frame
        set_text(tf2, subtitle, font_size=14, color=SSW_GRAY)


def new_slide(prs, slide_num, title, subtitle=None):
    """Create a new blank slide with section title and slide number."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, title, subtitle)
    add_slide_number(slide, slide_num)
    return slide


def content_box(slide, top=1.8, height=5.0, left=0.8, width=11.5):
    """Add a word-wrapped text box for body content."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    set_text(tf, "", font_size=12)  # clear
    return tf


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def create_title_slide(prs, data):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Red accent bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = SSW_RED
    shape.line.fill.background()

    # SSW squares motif
    for i, color in enumerate([SSW_RED, SSW_CHARCOAL, SSW_GRAY, RGBColor(0xE2, 0x52, 0x52)]):
        sq = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8 + i * 0.35), Inches(2.5),
            Inches(0.25), Inches(0.25)
        )
        sq.fill.solid()
        sq.fill.fore_color.rgb = color
        sq.line.fill.background()

    # Survey name
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(11.5), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    set_text(tf, data["metadata"]["surveyName"], font_size=36, bold=True, color=SSW_CHARCOAL)

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.5))
    tf2 = txBox2.text_frame
    set_text(tf2, "Survey Analysis Report", font_size=20, color=SSW_GRAY)

    # Date and response count
    meta = data["metadata"]
    date_str = meta.get("dateRange", "")
    count = meta.get("responseCount", "")
    info_text = f"{date_str}  •  {count} responses  •  Compulsory survey"

    txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.5))
    tf3 = txBox3.text_frame
    set_text(tf3, info_text, font_size=14, color=SSW_GRAY)


def create_executive_summary_slide(prs, data, slide_num):
    """Slide 2: Executive Summary with verdict."""
    slide = new_slide(prs, slide_num, "Executive Summary")

    # Resolve paths: overview.executiveSummary or top-level executiveSummary
    exec_summary = resolve(data, "overview.executiveSummary", "executiveSummary", default={})
    verdict = exec_summary.get("overallVerdict", "")

    # Verdict banner
    if verdict:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = SSW_CHARCOAL
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        set_text(tf, f"Verdict: {verdict}", font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Bullets
    bullets = exec_summary.get("bullets", [])
    tf = content_box(slide, top=2.8, height=4.0)
    write_bullets(tf, bullets[:5], font_size=16)


def create_key_metrics_slide(prs, data, slide_num):
    """Slide 3: Key Metrics."""
    slide = new_slide(prs, slide_num, "Key Metrics")

    metrics = data.get("keyMetrics", [])
    cols = min(len(metrics), 4)
    if cols == 0:
        return

    card_width = Inches(2.6)
    gap = Inches(0.3)
    total_width = cols * card_width + (cols - 1) * gap
    start_x = (SLIDE_WIDTH - total_width) / 2

    for i, metric in enumerate(metrics[:4]):
        x = start_x + i * (card_width + gap)
        y = Inches(2.0)

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, Inches(2.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = SSW_LIGHT_GRAY
        card.line.fill.background()

        # Metric value
        txVal = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.3), card_width - Inches(0.4), Inches(1.0))
        tf = txVal.text_frame
        tf.word_wrap = True
        set_text(tf, str(metric.get("value", "")), font_size=32, bold=True, color=SSW_RED, alignment=PP_ALIGN.CENTER)

        # Metric label
        txLbl = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.3), card_width - Inches(0.4), Inches(0.5))
        tf2 = txLbl.text_frame
        tf2.word_wrap = True
        set_text(tf2, metric.get("label", ""), font_size=13, bold=True, color=SSW_CHARCOAL, alignment=PP_ALIGN.CENTER)

        # Context
        context = metric.get("context", metric.get("status", ""))
        if context:
            txCtx = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.8), card_width - Inches(0.4), Inches(0.5))
            tf3 = txCtx.text_frame
            tf3.word_wrap = True
            set_text(tf3, str(context), font_size=10, color=SSW_GRAY, alignment=PP_ALIGN.CENTER)


def create_focus_area_slide(prs, data, slide_num):
    """Slide (optional): Focus Area Summary."""
    focus = resolve(data, "overview.focusSummary", "focusSummary", default=None)
    if not focus:
        return None

    focus_area = data.get("metadata", {}).get("focusArea", "Focus Area")
    slide = new_slide(prs, slide_num, f"Focus Area: {focus_area}")

    tf = content_box(slide)

    if isinstance(focus, str):
        set_text(tf, focus, font_size=16, color=SSW_CHARCOAL)
    elif isinstance(focus, dict):
        set_text(tf, focus.get("summary", str(focus)), font_size=16, color=SSW_CHARCOAL)
    elif isinstance(focus, list):
        write_bullets(tf, focus, font_size=16)

    return slide


def create_scores_slide(prs, data, slide_num):
    """Top & Bottom Scores (may produce 2 slides if many questions)."""
    slide = new_slide(prs, slide_num, "Top & Bottom Scores", "Highest and lowest scoring questions")

    questions = resolve(data, "responses.questionBreakdown", "questionBreakdown", default=[])
    if not questions:
        return 0

    sorted_q = sorted(questions, key=lambda q: q.get("mean", 0))
    bottom = sorted_q[:5]
    top = sorted(questions, key=lambda q: q.get("mean", 0), reverse=True)[:5]

    # Bottom scores (left column)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.4))
    set_text(txBox.text_frame, "⚠️ Lowest Scores", font_size=18, bold=True, color=SSW_RED)

    tf_bottom = content_box(slide, top=2.3, height=4.5, left=0.8, width=5.5)
    lines = []
    for q in bottom:
        text = q.get("text", "")[:70]
        mean = q.get("mean", 0)
        flags = q.get("flags", [])
        flag_str = f" [{', '.join(flags[:2])}]" if flags else ""
        lines.append(f"{text} — {mean:.1f}{flag_str}")
    write_bullets(tf_bottom, lines, font_size=13)

    # Top scores (right column)
    txBox2 = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.5), Inches(0.4))
    set_text(txBox2.text_frame, "✅ Highest Scores", font_size=18, bold=True, color=SSW_GREEN)

    tf_top = content_box(slide, top=2.3, height=4.5, left=7.0, width=5.5)
    lines = []
    for q in top:
        text = q.get("text", "")[:70]
        mean = q.get("mean", 0)
        lines.append(f"{text} — {mean:.1f}")
    write_bullets(tf_top, lines, font_size=13)

    return 1


def create_all_scores_slide(prs, data, slide_num):
    """Full question breakdown — all questions ranked by score."""
    questions = resolve(data, "responses.questionBreakdown", "questionBreakdown", default=[])
    if not questions or len(questions) <= 10:
        return 0

    sorted_q = sorted(questions, key=lambda q: q.get("mean", 0))
    slides_created = 0
    chunk_size = 12

    for chunk_start in range(0, len(sorted_q), chunk_size):
        chunk = sorted_q[chunk_start:chunk_start + chunk_size]
        page = chunk_start // chunk_size + 1
        total_pages = (len(sorted_q) + chunk_size - 1) // chunk_size
        subtitle = f"All questions ranked by score ({page}/{total_pages})"
        slide = new_slide(prs, slide_num + slides_created, "Question Scores", subtitle)

        tf = content_box(slide, top=1.6, height=5.5)
        lines = []
        for q in chunk:
            text = q.get("text", "")[:65]
            mean = q.get("mean", 0)
            flags = q.get("flags", [])
            commentary = q.get("commentary", "")
            flag_str = f" [{', '.join(flags[:2])}]" if flags else ""
            lines.append(f"{text} — {mean:.1f}{flag_str}")
            if commentary:
                lines.append(f"   {commentary[:120]}")
        write_bullets(tf, lines, font_size=11, prefix="")
        slides_created += 1

    return slides_created


def create_themes_slide(prs, data, slide_num):
    """Themes & Sentiment — may produce multiple slides."""
    slide = new_slide(prs, slide_num, "Themes & Sentiment", "What people are saying in their own words")

    # Sentiment overview banner
    sentiment = data.get("sentimentOverview", {})
    spectrum_label = sentiment.get("spectrumLabel", "")
    dominant = sentiment.get("dominantEmotion", "")
    key_insight = sentiment.get("keyInsight", "")

    banner_text = f"Emotional Temperature: {spectrum_label}"
    if dominant:
        banner_text += f"  |  Dominant: {dominant}"

    if spectrum_label:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.6)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = SSW_CHARCOAL
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        set_text(tf, banner_text, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Key insight
    start_y = 2.5
    if key_insight:
        txIns = slide.shapes.add_textbox(Inches(0.8), Inches(start_y), Inches(11.5), Inches(0.6))
        tf_ins = txIns.text_frame
        tf_ins.word_wrap = True
        set_text(tf_ins, key_insight, font_size=13, color=SSW_GRAY, alignment=PP_ALIGN.LEFT)
        start_y = 3.2

    # Top themes
    themes_data = data.get("themes", {})
    themes = themes_data.get("themes", []) if isinstance(themes_data, dict) else themes_data

    tf = content_box(slide, top=start_y, height=7.5 - start_y - 0.5)
    lines = []
    for theme in themes[:6]:
        name = theme.get("name", "")
        freq = theme.get("frequency", "")
        sentiment_val = theme.get("sentiment", "")
        intensity = theme.get("intensity", "")
        rep_quote = theme.get("representativeQuote", {})
        quote_text = rep_quote.get("text", "") if isinstance(rep_quote, dict) else str(rep_quote)
        respondent = rep_quote.get("name", rep_quote.get("respondent", "")) if isinstance(rep_quote, dict) else ""

        header = f"{name} — {freq} mentions, {sentiment_val}"
        if intensity:
            header += f" ({intensity} intensity)"
        lines.append(header)
        if quote_text:
            lines.append(f'   "{quote_text[:120]}" — {respondent}')

    write_bullets(tf, lines, font_size=13, prefix="")

    # If there are many themes, add a second slide
    slides_created = 1
    if len(themes) > 6:
        slide2 = new_slide(prs, slide_num + 1, "Themes (continued)")
        tf2 = content_box(slide2, top=1.6, height=5.5)
        lines = []
        for theme in themes[6:12]:
            name = theme.get("name", "")
            freq = theme.get("frequency", "")
            sentiment_val = theme.get("sentiment", "")
            rep_quote = theme.get("representativeQuote", {})
            quote_text = rep_quote.get("text", "") if isinstance(rep_quote, dict) else str(rep_quote)
            respondent = rep_quote.get("name", rep_quote.get("respondent", "")) if isinstance(rep_quote, dict) else ""

            lines.append(f"{name} — {freq} mentions, {sentiment_val}")
            if quote_text:
                lines.append(f'   "{quote_text[:120]}" — {respondent}')
        write_bullets(tf2, lines, font_size=13, prefix="")
        slides_created += 1

    return slides_created


def create_notable_quotes_slide(prs, data, slide_num):
    """Notable quotes from across the survey."""
    themes_data = data.get("themes", {})
    quotes = themes_data.get("notableQuotes", []) if isinstance(themes_data, dict) else []
    if not quotes:
        return 0

    slide = new_slide(prs, slide_num, "Notable Quotes", "Voices worth hearing")

    tf = content_box(slide, top=1.6, height=5.5)
    lines = []
    for q in quotes[:8]:
        name = q.get("name", "")
        text = q.get("text", "")[:150]
        question = q.get("question", "")[:80]
        category = q.get("category", "")
        why = q.get("whySelected", "")

        tag = f" [{category}]" if category else ""
        lines.append(f'"{text}"')
        attr = f"   — {name}"
        if question:
            attr += f', answering "{question}"'
        if why:
            attr += f" | {why[:80]}"
        lines.append(attr + tag)

    write_bullets(tf, lines, font_size=13, prefix="")
    return 1


def create_standout_responses_slide(prs, data, slide_num):
    """Standout responses — individual answers worth highlighting."""
    standouts = resolve(data, "overview.standoutResponses", "standoutResponses", default=[])
    if not standouts:
        return 0

    slide = new_slide(prs, slide_num, "Standout Responses", "Individual answers worth highlighting")

    tf = content_box(slide, top=1.6, height=5.5)
    lines = []
    for s in standouts[:6]:
        name = s.get("name", s.get("respondent", ""))
        question = s.get("question", "")[:80]
        response = s.get("response", "")[:150]
        why = s.get("whyStandout", "")

        lines.append(f"{name}")
        if question:
            lines.append(f'   Q: "{question}"')
        lines.append(f'   "{response}"')
        if why:
            lines.append(f"   Why: {why[:100]}")

    write_bullets(tf, lines, font_size=12, prefix="")
    return 1


def create_red_flags_slides(prs, data, slide_num):
    """Red Flags — may produce multiple slides for many flags."""
    flags = data.get("redFlags", [])
    if not flags:
        slide = new_slide(prs, slide_num, "🚩 Red Flags")
        tf = content_box(slide)
        set_text(tf, "No critical red flags identified.", font_size=16, color=SSW_GREEN)
        return 1

    slides_created = 0
    chunk_size = 4

    for chunk_start in range(0, len(flags), chunk_size):
        chunk = flags[chunk_start:chunk_start + chunk_size]
        page_suffix = f" ({chunk_start // chunk_size + 1}/{(len(flags) + chunk_size - 1) // chunk_size})" if len(flags) > chunk_size else ""

        slide = new_slide(prs, slide_num + slides_created,
                          f"🚩 Red Flags{page_suffix}",
                          "Critical warnings requiring attention")

        tf = content_box(slide, top=1.6, height=5.5)
        lines = []
        for flag in chunk:
            flag_text = flag.get("flag", "")
            severity = flag.get("severity", "")
            evidence = flag.get("evidence", "")
            prediction = flag.get("prediction", "")
            time_to_act = flag.get("timeToAct", "")

            lines.append(f"⚠️ {flag_text} [{severity.upper()}]")
            if evidence:
                lines.append(f"   Evidence: {evidence[:180]}")
            if prediction:
                lines.append(f"   Prediction: {prediction[:150]}")
            if time_to_act:
                lines.append(f"   Time to act: {time_to_act}")

        write_bullets(tf, lines, font_size=12, prefix="")
        slides_created += 1

    return slides_created


def create_recommendations_slides(prs, data, slide_num):
    """Recommendations — one slide per tier for full detail."""
    recs = data.get("recommendations", {})
    slides_created = 0

    tiers = [
        ("🔴 Immediate Actions (This Week)", recs.get("immediate", []), SSW_RED),
        ("🟡 Short-Term Actions (This Quarter)", recs.get("shortTerm", []), SSW_AMBER),
        ("🟢 Strategic Actions (This Year)", recs.get("strategic", []), SSW_GREEN),
    ]

    for tier_name, items, color in tiers:
        if not items:
            continue

        slide = new_slide(prs, slide_num + slides_created, tier_name)

        tf = content_box(slide, top=1.6, height=5.5)
        lines = []
        for item in items:
            if isinstance(item, dict):
                action = item.get("action", "")
                owner = item.get("owner", "")
                rationale = item.get("rationale", "")
                metric = item.get("successMetric", "")

                lines.append(f"➡️ {action}")
                if owner:
                    lines.append(f"   Owner: {owner}")
                if rationale:
                    lines.append(f"   Why: {rationale[:150]}")
                if metric:
                    lines.append(f"   Success metric: {metric[:120]}")
            else:
                lines.append(f"➡️ {str(item)}")

        write_bullets(tf, lines, font_size=12, prefix="")
        slides_created += 1

    return slides_created


def create_emotional_profile_slide(prs, data, slide_num):
    """Emotional profile breakdown."""
    sentiment = data.get("sentimentOverview", {})
    breakdown = sentiment.get("emotionalBreakdown", {})
    if not breakdown:
        return 0

    slide = new_slide(prs, slide_num, "Emotional Profile", "How people feel across the survey")

    # Candor info
    candor = sentiment.get("candorLevel", "")
    dissonance = sentiment.get("quantQualDissonance", "")
    detail = sentiment.get("quantQualDetail", {})
    adjusted = detail.get("adjustedScoreEstimate", "") if isinstance(detail, dict) else ""

    info_parts = []
    if candor:
        info_parts.append(f"Candor level: {candor}")
    if dissonance:
        info_parts.append(f"Quant-qual dissonance: {dissonance}%")
    if adjusted:
        info_parts.append(f"Adjusted score estimate: {adjusted}")

    if info_parts:
        txInfo = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.4))
        set_text(txInfo.text_frame, "  |  ".join(info_parts), font_size=12, color=SSW_GRAY)

    # Emotional bars
    sorted_emotions = sorted(breakdown.items(), key=lambda x: x[1], reverse=True)

    y_start = 2.2
    bar_height = 0.35
    max_val = max(breakdown.values()) if breakdown.values() else 1
    max_bar_width = 7.0

    for i, (emotion, value) in enumerate(sorted_emotions):
        y = Inches(y_start + i * (bar_height + 0.15))

        # Label
        txLbl = slide.shapes.add_textbox(Inches(0.8), y, Inches(2.5), Inches(bar_height))
        tf = txLbl.text_frame
        set_text(tf, f"{emotion.capitalize()}", font_size=12, bold=True, color=SSW_CHARCOAL, alignment=PP_ALIGN.RIGHT)

        # Bar
        bar_width = max(0.1, (value / max_val) * max_bar_width) if max_val > 0 else 0.1
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(3.5), y + Inches(0.05),
            Inches(bar_width), Inches(bar_height - 0.1)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = SSW_RED
        bar.line.fill.background()

        # Value
        txVal = slide.shapes.add_textbox(Inches(3.5 + bar_width + 0.15), y, Inches(1.0), Inches(bar_height))
        set_text(txVal.text_frame, f"{value}%", font_size=11, color=SSW_GRAY)

    return 1


def create_hard_truths_slide(prs, data, slide_num):
    """Hard Truths."""
    truths = data.get("hardTruths", [])
    if not truths:
        return 0

    slide = new_slide(prs, slide_num, "Hard Truths", "What leadership needs to hear")

    # Red accent card
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(2.0), Inches(10.3), Inches(4.0)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = SSW_RED_50
    shape.line.color.rgb = SSW_RED
    shape.line.width = Pt(2)

    txBox = slide.shapes.add_textbox(Inches(2.0), Inches(2.5), Inches(9.3), Inches(3.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, truth in enumerate(truths[:2]):
        if i == 0:
            set_text(tf, truth, font_size=18, color=SSW_CHARCOAL)
        else:
            add_paragraph(tf, truth, font_size=18, color=SSW_CHARCOAL, space_before=20)

    return 1


def create_cross_survey_slide(prs, data, slide_num):
    """Cross-survey synthesis (multi-survey only)."""
    synthesis = data.get("crossSurveySynthesis", None)
    if not synthesis:
        return 0

    slide = new_slide(prs, slide_num, "Cross-Survey Patterns", "Where multiple surveys tell the same story")

    tf = content_box(slide, top=1.6, height=5.5)
    lines = []
    for item in synthesis[:6]:
        if isinstance(item, dict):
            topic = item.get("topic", "")
            surveys = item.get("surveys", [])
            evidence = item.get("evidence", "")
            insight = item.get("insight", "")
            strength = item.get("strength", "")

            tag = f" [{strength}]" if strength else ""
            lines.append(f"{topic}{tag}")
            if surveys:
                lines.append(f"   Surveys: {', '.join(surveys)}")
            if evidence:
                lines.append(f"   Evidence: {evidence[:150]}")
            if insight:
                lines.append(f"   {insight[:150]}")
        else:
            lines.append(str(item))

    write_bullets(tf, lines, font_size=12, prefix="")
    return 1


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    if len(sys.argv) != 3:
        print(f"Usage: {sys.argv[0]} <consolidated.json> <output.pptx>")
        sys.exit(1)

    input_path = Path(sys.argv[1])
    output_path = Path(sys.argv[2])

    if not input_path.exists():
        print(f"Error: {input_path} not found")
        sys.exit(1)

    with open(input_path, "r") as f:
        data = json.load(f)

    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 13.333 inches
    prs.slide_height = Emu(6858000)   # 7.5 inches

    slide_num = 1

    # 1. Title
    create_title_slide(prs, data)
    slide_num += 1

    # 2. Executive Summary + Verdict
    create_executive_summary_slide(prs, data, slide_num)
    slide_num += 1

    # 3. Key Metrics
    create_key_metrics_slide(prs, data, slide_num)
    slide_num += 1

    # 4. Focus Area (optional)
    if create_focus_area_slide(prs, data, slide_num):
        slide_num += 1

    # 5. Top & Bottom Scores
    n = create_scores_slide(prs, data, slide_num)
    slide_num += max(n, 1)

    # 6+. All Question Scores (paginated)
    n = create_all_scores_slide(prs, data, slide_num)
    slide_num += n

    # 7. Themes & Sentiment (may be 1-2 slides)
    n = create_themes_slide(prs, data, slide_num)
    slide_num += n

    # 8. Emotional Profile
    n = create_emotional_profile_slide(prs, data, slide_num)
    slide_num += n

    # 9. Notable Quotes
    n = create_notable_quotes_slide(prs, data, slide_num)
    slide_num += n

    # 10. Standout Responses
    n = create_standout_responses_slide(prs, data, slide_num)
    slide_num += n

    # 11. Cross-Survey Patterns (optional)
    n = create_cross_survey_slide(prs, data, slide_num)
    slide_num += n

    # 12+. Red Flags (paginated)
    n = create_red_flags_slides(prs, data, slide_num)
    slide_num += n

    # 13+. Recommendations (one slide per tier)
    n = create_recommendations_slides(prs, data, slide_num)
    slide_num += n

    # 14. Hard Truths (optional)
    n = create_hard_truths_slide(prs, data, slide_num)
    slide_num += n

    # Save
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Slide deck saved to {output_path} ({slide_num - 1} slides)")


if __name__ == "__main__":
    main()
